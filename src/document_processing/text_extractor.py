import pdfplumber
import docx
import openpyxl
from pptx import Presentation
import os
import io

def _extract_text_from_pdf(file_obj):
    full_text = []
    with pdfplumber.open(file_obj) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text: full_text.append(page_text)
    return "\n".join(full_text)

def _extract_text_from_docx(file_obj):
    doc = docx.Document(file_obj)
    return "\n".join([para.text for para in doc.paragraphs])

def _extract_text_from_excel(file_obj):
    workbook = openpyxl.load_workbook(file_obj)
    full_text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows():
            row_text = [str(cell.value) for cell in row if cell.value is not None]
            full_text.append(", ".join(row_text))
    return "\n".join(full_text)
    
def _extract_text_from_pptx(file_obj):
    prs = Presentation(file_obj)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"): full_text.append(shape.text)
    return "\n".join(full_text)

def extract_text(file_source, file_extension=None):
    """
    Extracts text from a file source, which can be a file path (string)
    or a file-like object (bytes).
    """
    if isinstance(file_source, str): 
        if not file_extension:
            _, file_extension = os.path.splitext(file_source)
        with open(file_source, 'rb') as f:
            file_obj = io.BytesIO(f.read())
    else: 
        file_obj = file_source

    ext = file_extension.lower()
    if ext == '.pdf': 
        return _extract_text_from_pdf(file_obj)
    elif ext == '.docx': 
        return _extract_text_from_docx(file_obj)
    elif ext == '.xlsx': 
        return _extract_text_from_excel(file_obj)
    elif ext == '.pptx': 
        return _extract_text_from_pptx(file_obj)
    else: 
        return f"Unsupported file type: {ext}"